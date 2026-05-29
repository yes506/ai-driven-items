#!/usr/bin/env python3
"""validate_anchors.py — cross-reference resolution for text + structured output.

Two modes:

  --text   TARGET_PATH is markdown; every Markdown `[…](#anchor)` link
           must resolve to a heading or `<a id=…>` anchor in the same
           file. Literal `[[stub-id]]` wikilinks in TARGET_PATH are
           reported as un-transformed planner-internal syntax.

  --pptx   TARGET_PATH is .pptx. Opens with python-pptx (non-mutating),
           verifies the file parses + slide count > 0. When --plan is
           provided, additionally verifies slide count >= declared-stub
           count and each stub-id appears as the first line of some
           slide's speaker notes (provenance contract).

`--plan <document-plan.md>` is OPTIONAL:
  - feature/system: pass it. Enables completeness/provenance checks
    against declared stubs.
  - micro/local: omit it. The implementer's bullet-derived work queue
    has no document-plan.md artifact (chat-only planner contract).
    Validators fall back to standalone target checks.

Usage:
  validate_anchors.py --text [--plan <document-plan.md>] <target.md>
  validate_anchors.py --pptx [--plan <document-plan.md>] <target.pptx>

Exits 0 on success, 1 on validation failure, 2 on usage / dep error.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

STUB_HEADING = re.compile(r"^\s*##\s+stub:\s*([A-Za-z0-9_-]+)\s*$")
WIKILINK = re.compile(r"\[\[([A-Za-z0-9_-]+)\]\]")
MD_ANCHOR = re.compile(r"\]\(#([A-Za-z0-9_-]+)\)")


def slugify_heading(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"\s+", "-", text)
    return text.strip("-")


def declared_stub_ids(plan_path: Path) -> list[str]:
    return [
        match.group(1)
        for raw in plan_path.read_text(encoding="utf-8").splitlines()
        if (match := STUB_HEADING.match(raw))
    ]


def validate_text(plan_path: Path | None, target_path: Path) -> list[str]:
    errors: list[str] = []
    text = target_path.read_text(encoding="utf-8")

    anchors: set[str] = set()
    for raw in text.splitlines():
        heading = re.match(r"^\s*#{1,6}\s+(.+?)\s*$", raw)
        if heading:
            anchors.add(slugify_heading(heading.group(1)))
    for stub_id in re.findall(r"<a\s+id=[\"']([A-Za-z0-9_-]+)[\"']", text):
        anchors.add(stub_id)

    for lineno, raw in enumerate(text.splitlines(), start=1):
        for anchor in MD_ANCHOR.findall(raw):
            if anchor not in anchors:
                errors.append(
                    f"line {lineno}: markdown anchor `#{anchor}` does not resolve"
                )
        for wiki in WIKILINK.findall(raw):
            errors.append(
                f"line {lineno}: untransformed wikilink `[[{wiki}]]` "
                "(implementer must transform to target-format anchor)"
            )

    return errors


def validate_pptx(
    plan_path: Path | None, target_path: Path
) -> tuple[int, list[str]]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        return 2, [
            "python-pptx not installed; "
            "install with `pip install python-pptx` or skip --pptx validation"
        ]

    errors: list[str] = []

    try:
        deck = Presentation(str(target_path))
    except Exception as e:
        return 1, [f"failed to open {target_path}: {e}"]

    slides = list(deck.slides)
    if not slides:
        return 1, [f"{target_path}: 0 slides — empty deck"]

    if plan_path is None:
        # Micro/local mode — no document-plan.md exists. Just verify the
        # file parses and has at least one slide. Skip provenance checks.
        return 0, errors

    # Feature/system mode — full provenance check.
    declared = declared_stub_ids(plan_path)
    declared_set = set(declared)
    if len(slides) < len(declared):
        errors.append(
            f"slide count ({len(slides)}) is less than declared stub count "
            f"({len(declared)})"
        )

    seen_stub_ids: set[str] = set()
    for idx, slide in enumerate(slides, start=1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        first_line = notes_text.splitlines()[0] if notes_text else ""
        if first_line in declared_set:
            seen_stub_ids.add(first_line)
        else:
            errors.append(
                f"slide {idx}: speaker-notes first line `{first_line!r}` "
                "is not a declared stub id (provenance contract violation)"
            )

    missing = sorted(declared_set - seen_stub_ids)
    for stub_id in missing:
        errors.append(f"declared stub `{stub_id}` has no provenance slide")

    return (1 if errors else 0), errors


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(add_help=False)
    parser.add_argument("--text", action="store_true")
    parser.add_argument("--pptx", action="store_true")
    parser.add_argument(
        "--plan",
        dest="plan_path",
        default=None,
        help="document-plan.md (feature/system only; omit for micro/local)",
    )
    parser.add_argument("target_path")
    try:
        args = parser.parse_args(argv[1:])
    except SystemExit:
        sys.stderr.write(
            f"usage: {argv[0]} (--text | --pptx) "
            "[--plan <document-plan.md>] <target>\n"
        )
        return 2

    if args.text == args.pptx:  # both or neither
        sys.stderr.write("specify exactly one of --text or --pptx\n")
        return 2

    plan_path: Path | None = None
    if args.plan_path is not None:
        plan_path = Path(args.plan_path)
        if not plan_path.is_file():
            sys.stderr.write(f"plan file not found: {plan_path}\n")
            return 2

    target_path = Path(args.target_path)
    if not target_path.is_file():
        sys.stderr.write(f"target file not found: {target_path}\n")
        return 2

    if args.text:
        errors = validate_text(plan_path, target_path)
        exit_code = 1 if errors else 0
    else:
        exit_code, errors = validate_pptx(plan_path, target_path)

    if errors:
        sys.stderr.write(f"validate_anchors: FAIL ({len(errors)} issue(s))\n")
        for e in errors:
            sys.stderr.write(f"  - {e}\n")
        return exit_code

    sys.stdout.write(f"validate_anchors: OK ({target_path.name})\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
