#!/usr/bin/env python3
"""render_pptx.py — emit .pptx from the implementer state's work queue.

For OUTPUT_STACK=structured (DOCTYPE=ppt) only. Reads each completed item
in `.document-implementer-state.json::work_queue[]` (the canonical state
shape per `references/state-and-resume.md`):

  - stub_id = item.item_id (or item.spec_payload.id as fallback)
  - title   = item.spec_payload.title (or item.spec_payload.purpose, or stub_id)
  - body    = item.generated_content

Each slide carries its source stub_id as the first line of speaker notes
— `validate_anchors.py --pptx` reads this back as the provenance contract.

Usage:
  render_pptx.py <state.json> <output.pptx>

Exits 0 on success, 1 on render failure, 2 on usage error or missing dep.

Imports stdlib + python-pptx. No PyYAML. python-pptx availability checked
by the implementer's Phase 1 dep-check; absence here is a fatal error
because the implementer should never have invoked this script otherwise.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path


def main(argv: list[str]) -> int:
    if len(argv) != 3:
        sys.stderr.write(f"usage: {argv[0]} <state.json> <output.pptx>\n")
        return 2

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.util import Inches, Pt  # type: ignore[import-not-found]
    except ImportError:
        sys.stderr.write(
            "python-pptx not installed; "
            "the implementer's Phase 1 dep-check should have prevented this. "
            "Install with `pip install python-pptx` and re-run.\n"
        )
        return 2

    state_path = Path(argv[1])
    out_path = Path(argv[2])
    if not state_path.is_file():
        sys.stderr.write(f"not a file: {state_path}\n")
        return 2

    try:
        state = json.loads(state_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as e:
        sys.stderr.write(f"cannot parse {state_path}: {e}\n")
        return 2

    work_queue = state.get("work_queue", [])
    completed = [
        item for item in work_queue if item.get("status") == "completed"
    ]
    if not completed:
        sys.stderr.write(
            "state.work_queue has no completed items to render "
            f"(found {len(work_queue)} total, "
            f"{sum(1 for i in work_queue if i.get('status') == 'completed')} completed)\n"
        )
        return 1

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # "Title Only" — title + body box

    for item in completed:
        spec = item.get("spec_payload", {}) or {}
        stub_id = item.get("item_id") or spec.get("id") or ""
        title = spec.get("title") or spec.get("purpose") or stub_id
        body = item.get("generated_content", "")

        slide = prs.slides.add_slide(blank_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        first = True
        for line in (body or "").splitlines() or [""]:
            paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(18)
            first = False

        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = stub_id

    try:
        prs.save(str(out_path))
    except OSError as e:
        sys.stderr.write(f"failed to write {out_path}: {e}\n")
        return 1

    sys.stdout.write(
        f"render_pptx: OK ({len(completed)} slide(s) → {out_path})\n"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
